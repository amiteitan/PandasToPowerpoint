import six
from pptx import Presentation
from pptx.util import Cm, Pt
import pandas as pd
from math import *

round_to_n = lambda x, n: round(x, -int(floor(log10(abs(x)))) + (n - 1))

def _do_formatting(value, format_str):
    """Format value according to format_str, and deal
    sensibly with format_str if it is missing or invalid."""
    if format_str == '':
        if type(value) in six.integer_types:
            format_str = ','
        elif type(value) is float:
            format_str = 'f'
        elif type(value) is str:
            format_str = 's'
    elif format_str[0] == '.':
        if format_str.endswith('R'):
            if type(value) in six.integer_types:
                value = round_to_n(value, int(format_str[1]))
                format_str = ','
        if not format_str.endswith('G'):
            format_str = format_str + "G"
    try:
        value = format(value, format_str)
    except:
        value = format(value, '')

    return value

def process_position_parameter(param):
    """Process positioning parameters (left, top, width, height) given to
    df_to_table.

    If an integer, returns the right instance of the Cm class to allow it to be treated
    as cm. If missing, then default to 4cm. Otherwise, pass through whatever it gets.
    """
    if param is None:
        return Cm(4)
    elif type(param) is int:
        return Cm(param)
    else:
        return param


def df_to_table(slide, df, left=None, top=None, width=None, height=None,
                colnames=None, col_formatters=None, rounding=None):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.
    
    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.
    
    Arguments:
     - slide: slide object from the python-pptx library containing the slide on which you want the table to appear
     - df: Pandas DataFrame with the data
     
    Optional arguments:
     - left: Position of the left-side of the table, either as an integer in cm, or as an instance of a
     pptx.util Length class (pptx.util.Inches for example). Defaults to 4cm.
     - top: Position of the top of the table, takes parameters as above.
     - width: Width of the table, takes parameters as above.
     - height: Height of the table, takes parameters as above.
     - col_formatters: A n_columns element long list containing format specifications for each column.
     For example ['', ',', '.2'] does no special formatting for the first column, uses commas as thousands separators
     in the second column, and formats the third column as a float with 2 decimal places.
     - rounding: A n_columns element long list containing a number for each integer column that requires rounding
     that is then multiplied by -1 and passed to round(). The practical upshot of this is that you can give something like
     ['', 3, ''], which does nothing for the 1st and 3rd columns (as they aren't integer values), but for the 2nd column,
     rounds away the 3 right-hand digits (eg. taking 25437 to 25000).
     """
    left = process_position_parameter(left)
    top = process_position_parameter(top)
    width = process_position_parameter(width)
    height = process_position_parameter(height)

    rows, cols = df.shape
    res = slide.shapes.add_table(rows+1, cols, left, top, width, height)
    
    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        res.table.cell(0,col_index).text = col_name
        
    m = df.as_matrix()
    
    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            
            if col_formatters is None:
                text = str(val)
            else:
                #text = col_formatters[col].format(m[row, col])
                text = _do_formatting(val, col_formatters[col])
            
            res.table.cell(row+1, col).text = text
            #res.table.cell(row+1, col).text_frame.fit_text()
    
def df_to_powerpoint(filename, df, **kwargs):
    """Converts a Pandas DataFrame to a table in a new, blank PowerPoint presentation.
    
    Creates a new PowerPoint presentation with the given filename, with a single slide
    containing a single table with the Pandas DataFrame data in it.
    
    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.
    
    Arguments:
     - filename: Filename to save the PowerPoint presentation as
     - df: Pandas DataFrame with the data
    
    All other arguments that can be taken by df_to_table() (such as col_formatters or rounding) can also
    be passed here.
    """
    pres = Presentation()
    blank_slide_layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide_layout)
    df_to_table(slide, df, **kwargs)
    pres.save(filename)
